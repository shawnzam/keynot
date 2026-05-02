#!/usr/bin/env python3
"""
Proof-of-concept: recreate the keynot-for-zombies deck as a .pptx using python-pptx.
Run with: uv run --with python-pptx examples/keynot-for-zombies-pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors from the HTML deck
PRIMARY = RGBColor(0x1A, 0x2E, 0x1F)
ACCENT = RGBColor(0xA3, 0xE6, 0x35)
BLOOD = RGBColor(0x8B, 0x1A, 0x1A)
CREAM = RGBColor(0xF3, 0xEF, 0xE4)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x6A, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(10)

# Use blank layout
blank_layout = prs.slide_layouts[6]


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, italic=False, color=WHITE, align=PP_ALIGN.LEFT,
                 font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


# ===================== SLIDE 1 — Title =====================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, PRIMARY)

add_text_box(slide1, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4),
             "KEYNOT", font_size=11, bold=True, color=ACCENT, font_name="Arial")
add_text_box(slide1, Inches(0.8), Inches(0.9), Inches(3), Inches(0.4),
             "HTML SLIDE SKILL", font_size=10, color=RGBColor(0x88, 0x88, 0x88), font_name="Arial")

add_text_box(slide1, Inches(0.8), Inches(1.8), Inches(6), Inches(0.5),
             "A Field Guide · Edition One", font_size=10, bold=True, color=ACCENT, font_name="Arial")

add_text_box(slide1, Inches(0.8), Inches(2.5), Inches(10), Inches(4),
             "Keynot for\nZombies", font_size=96, color=WHITE, font_name="Georgia")

add_text_box(slide1, Inches(0.8), Inches(7.0), Inches(10), Inches(1.5),
             "Slide decks so simple, even a reanimated corpse with half a\nbrain can make one. No apps. No crashes. Just slides.",
             font_size=22, italic=True, color=RGBColor(0xB0, 0xB0, 0xB0), font_name="Georgia")

add_text_box(slide1, Inches(0.8), Inches(9.2), Inches(5), Inches(0.4),
             "Chapter 00 · Introduction", font_size=10, color=RGBColor(0x70, 0x70, 0x70), font_name="Arial")


# ===================== SLIDE 2 — The Problem =====================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, CREAM)

# Left dark panel
left_panel = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(8), Inches(10))
left_panel.fill.solid()
left_panel.fill.fore_color.rgb = PRIMARY
left_panel.line.fill.background()

add_text_box(slide2, Inches(0.8), Inches(2.5), Inches(6), Inches(0.5),
             "THE CONDITION", font_size=10, bold=True, color=ACCENT, font_name="Arial")
add_text_box(slide2, Inches(0.8), Inches(3.2), Inches(6.5), Inches(2.5),
             "Your brain is…\nnot what it used to be.", font_size=48, color=WHITE, font_name="Georgia")
add_text_box(slide2, Inches(0.8), Inches(6.0), Inches(6), Inches(1.5),
             "You need slides by morning. You cannot remember which app does what. Every click makes it worse.",
             font_size=15, color=RGBColor(0x99, 0x99, 0x99), font_name="Arial")

# Right side — pain points
add_text_box(slide2, Inches(8.5), Inches(2.5), Inches(6), Inches(0.4),
             "SYMPTOMS YOU MAY RECOGNIZE", font_size=10, bold=True, color=MUTED, font_name="Arial")

pains = [
    ("01", "PowerPoint asks twelve questions before you can type. Then it crashes."),
    ("02", "Keynote only opens on one computer. Guess which one isn't yours."),
    ("03", "Google Slides looks like Google Slides. You wanted it to look good."),
    ("04", "You just want to present something. Today. In a browser. That's all."),
]
for i, (num, text) in enumerate(pains):
    y = 3.2 + i * 1.5
    add_text_box(slide2, Inches(8.5), Inches(y), Inches(0.6), Inches(0.5),
                 num, font_size=24, color=BLOOD, font_name="Georgia")
    add_text_box(slide2, Inches(9.2), Inches(y), Inches(6), Inches(1.2),
                 text, font_size=15, color=INK, font_name="Arial")


# ===================== SLIDE 3 — How it works =====================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3, CREAM)

add_text_box(slide3, Inches(0.8), Inches(1.0), Inches(6), Inches(0.4),
             "THE CURE · THREE EASY STEPS", font_size=10, bold=True, color=BLOOD, font_name="Arial")
add_text_box(slide3, Inches(0.8), Inches(1.6), Inches(10), Inches(2),
             "You talk. Claude types.\nA deck appears.", font_size=48, color=INK, font_name="Georgia")

steps = [
    ("01", "Say what you want", "\"Make me a deck about X for audience Y.\" One sentence is enough. Keynot will ask follow-ups if it needs them."),
    ("02", "Claude writes the HTML", "A single self-contained file. Brand colors, fonts, layouts, reveals — all baked in. No build step, no dependencies."),
    ("03", "Open it in a browser", "Arrow keys to navigate. F for fullscreen. Swipe on mobile. That's the whole thing."),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.8 + i * 5.0
    add_text_box(slide3, Inches(x), Inches(4.5), Inches(1.5), Inches(1),
                 num, font_size=56, color=ACCENT, font_name="Georgia")
    add_text_box(slide3, Inches(x + 1.5), Inches(4.5), Inches(3.2), Inches(0.6),
                 title, font_size=20, bold=True, color=INK, font_name="Georgia")
    add_text_box(slide3, Inches(x + 1.5), Inches(5.3), Inches(3.2), Inches(2.5),
                 desc, font_size=13, color=MUTED, font_name="Arial")


# ===================== SLIDE 4 — What you get =====================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4, PRIMARY)

# Accent bar at top
top_bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Pt(4))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = ACCENT
top_bar.line.fill.background()

add_text_box(slide4, Inches(0.8), Inches(0.8), Inches(6), Inches(0.4),
             "INCLUDED IN EVERY DECK", font_size=10, bold=True, color=ACCENT, font_name="Arial")
add_text_box(slide4, Inches(0.8), Inches(1.4), Inches(10), Inches(1.5),
             "One file.\nEverything inside.", font_size=48, color=WHITE, font_name="Georgia")

features = [
    ("Self-contained HTML", "Fonts via CDN, images as base64. Email it, AirDrop it, stick it on a thumb drive."),
    ("Keyboard & swipe nav", "Arrow keys, spacebar, touch swipes. Dots show progress. Nav auto-hides."),
    ("Fullscreen with one key", "Press F. No \"Presenter Mode\" wizard. No second monitor setup ritual."),
    ("Brand-matched design", "Hand it a style guide and it reads the colors, fonts, layout language."),
    ("Animated reveals", "Staggered fade-ups per slide, built in. No PowerPoint transition menu required."),
    ("Iterable by chat", "\"Slide 3: swap the headline.\" Done. No hunting through a sidebar of 40 objects."),
]
for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 7.5
    y = 3.8 + row * 2.0
    add_text_box(slide4, Inches(x), Inches(y), Inches(6.5), Inches(0.5),
                 title, font_size=20, bold=True, color=WHITE, font_name="Georgia")
    add_text_box(slide4, Inches(x), Inches(y + 0.5), Inches(6.5), Inches(1.2),
                 desc, font_size=13, color=RGBColor(0xA0, 0xA0, 0xA0), font_name="Arial")


# ===================== SLIDE 5 — CTA =====================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5, PRIMARY)

add_text_box(slide5, Inches(0.8), Inches(1.0), Inches(5), Inches(0.4),
             "FINAL THOUGHT", font_size=10, bold=True, color=ACCENT, font_name="Arial")
add_text_box(slide5, Inches(0.8), Inches(2.5), Inches(8), Inches(4),
             "Now go\neat some\nslides.", font_size=96, color=WHITE, font_name="Georgia")
add_text_box(slide5, Inches(0.8), Inches(7.0), Inches(8), Inches(1.5),
             "Open Claude Code. Type \"make me a deck about…\"\nWatch the skill wake up. The rest is autopilot.",
             font_size=20, italic=True, color=RGBColor(0xB0, 0xB0, 0xB0), font_name="Georgia")

# "braaains" on the right
add_text_box(slide5, Inches(10), Inches(3.5), Inches(5.5), Inches(4),
             "braaains", font_size=120, italic=True,
             color=RGBColor(0x2A, 0x4A, 0x32), font_name="Georgia", align=PP_ALIGN.CENTER)

# Save
output_path = "examples/keynot-for-zombies.pptx"
prs.save(output_path)
print(f"✓ Saved: {output_path}")
